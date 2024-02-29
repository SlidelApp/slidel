import os
from io import BytesIO
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def convert_pdf_to_images(pdf_file, output_folder):
    # Create output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    # Open the PDF file
    pdf_document = fitz.open(pdf_file)
    
    # Convert each page of the PDF into an image
    for page_number in range(len(pdf_document)):
        page = pdf_document.load_page(page_number)
        image = page.get_pixmap()
        image_bytes = BytesIO(image.tobytes())
        image_pil = Image.open(image_bytes)
        # Resize image to reduce memory usage
        image_pil.thumbnail((1024, 1024))
        image_pil.save(os.path.join(output_folder, f'page_{page_number + 1}.jpg'))

def convert_pptx_to_images(pptx_file, output_folder):
    # Create output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)
    
    # Convert each slide of the presentation into an image
    for i, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # If the shape is a picture, save it as an image
                image = shape.image
                image_bytes = BytesIO(image.blob)
                image_pil = Image.open(image_bytes)
                # Resize image to reduce memory usage
                image_pil.thumbnail((1024, 1024))
                image_pil.save(os.path.join(output_folder, f'slide_{i + 1}.jpg'))
                break
        else:
            # If no picture shape found, create a blank slide
            slide_image = Image.new('RGB', (int(presentation.slide_width), int(presentation.slide_height)), (255, 255, 255))
            # Resize image to reduce memory usage
            slide_image.thumbnail((1024, 1024))
            slide_image.save(os.path.join(output_folder, f'slide_{i + 1}.jpg'))

def convert_file_to_images(input_file, output_folder):
    # Determine the file type based on extension
    file_extension = os.path.splitext(input_file)[1].lower()
    if file_extension == '.pdf':
        convert_pdf_to_images(input_file, output_folder)
        print(f'PDF converted to images. Images saved in {output_folder} folder.')
    elif file_extension == '.pptx':
        convert_pptx_to_images(input_file, output_folder)
        print(f'PPTX converted to images. Images saved in {output_folder} folder.')
    else:
        print('Unsupported file format. Only PDF and PPTX files are supported.')

if __name__ == "__main__":
    # Specify the input file and output folder
    pdf_file = 'D:\\slidel\\slidel\\ui.py\\input.pdf'
    pptx_file = 'D:\\slidel\\slidel\\ui.py\\input.pptx'
    output_folder_pdf = 'pdf_images'
    output_folder_pptx = 'pptx_images'
    
    # Convert PDF to images
    convert_pdf_to_images(pdf_file, output_folder_pdf)
    print(f'PDF converted to images. Images saved in {output_folder_pdf} folder.')
    
    # Convert PPTX to images
    convert_pptx_to_images(pptx_file, output_folder_pptx)
    print(f'PPTX converted to images. Images saved in {output_folder_pptx} folder.')
